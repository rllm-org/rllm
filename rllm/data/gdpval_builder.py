"""Materialize OpenAI GDPval as rLLM sandbox tasks.

The source dataset stores prompts plus Hugging Face paths for task inputs and
expert deliverables.  This builder stages the inputs in the task environment,
keeps expert files under ``tests/reference`` (uploaded only after the agent
finishes), writes the GDPval solver image, and emits an LLM-judge verifier.
"""

from __future__ import annotations

import hashlib
import json
import logging
import posixpath
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree

logger = logging.getLogger(__name__)

REPO_ID = "openai/gdpval"
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_EXTERNAL_OPEN_REL = re.compile(
    rb"<Relationship\b(?P<attrs>[^<>]*\bTargetMode=(?:\"External\"|'External')[^<>]*?)(?<!/)>",
    re.IGNORECASE,
)
_SUSPICIOUS_SETTINGS = re.compile(rb"(?:xmlns:ns\d+\s*=|</?ns\d+:)")


def _toml_escape(value: str) -> str:
    return value.replace("\\", "\\\\").replace('"""', '\\"\\"\\"')


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _relationship_source(name: str) -> str | None:
    if name == "_rels/.rels":
        return ""
    marker = "/_rels/"
    if marker not in name or not name.endswith(".rels"):
        return None
    parent, leaf = name.split(marker, 1)
    return posixpath.join(parent, leaf[: -len(".rels")])


def _validate_office_package(path: Path) -> dict:
    report: dict = {
        "zip_ok": False,
        "required_parts_missing": [],
        "xml_errors": [],
        "missing_relationship_targets": [],
        "valid": False,
    }
    try:
        with zipfile.ZipFile(path) as package:
            names = set(package.namelist())
            bad_member = package.testzip()
            if bad_member:
                report["xml_errors"].append(f"CRC failure: {bad_member}")
                return report
            report["zip_ok"] = True
            report["required_parts_missing"] = sorted({"[Content_Types].xml", "_rels/.rels"} - names)
            for name in sorted(names):
                if not (name.endswith(".xml") or name.endswith(".rels")):
                    continue
                try:
                    root = ElementTree.fromstring(package.read(name))
                except (ElementTree.ParseError, UnicodeError) as exc:
                    report["xml_errors"].append(f"{name}: {exc}")
                    continue
                if not name.endswith(".rels"):
                    continue
                source = _relationship_source(name)
                if source is None:
                    continue
                source_dir = posixpath.dirname(source)
                for relationship in root.findall(f"{{{_REL_NS}}}Relationship"):
                    if relationship.get("TargetMode", "").lower() == "external":
                        continue
                    target = relationship.get("Target")
                    if not target or target.startswith("/"):
                        continue
                    resolved = posixpath.normpath(posixpath.join(source_dir, target))
                    if resolved not in names:
                        report["missing_relationship_targets"].append(f"{name} -> {resolved}")
    except (OSError, zipfile.BadZipFile) as exc:
        report["xml_errors"].append(str(exc))
        return report

    report["valid"] = not (report["required_parts_missing"] or report["xml_errors"] or report["missing_relationship_targets"])
    return report


def _repair_docx_members(package: zipfile.ZipFile) -> tuple[list[tuple[zipfile.ZipInfo, bytes]], list[str]]:
    members = [(info, package.read(info.filename)) for info in package.infolist()]
    by_name = {info.filename: data for info, data in members}
    changes: list[str] = []
    remove_settings = bool(_SUSPICIOUS_SETTINGS.search(by_name.get("word/settings.xml", b"")))
    if remove_settings:
        changes.append("removed suspicious word/settings.xml")

    repaired: list[tuple[zipfile.ZipInfo, bytes]] = []
    for info, data in members:
        if remove_settings and info.filename == "word/settings.xml":
            continue
        if info.filename.endswith(".rels"):
            data, count = _EXTERNAL_OPEN_REL.subn(rb"<Relationship\g<attrs>/>", data)
            if count:
                changes.append(f"closed {count} malformed external relationship(s) in {info.filename}")
            if remove_settings and info.filename == "word/_rels/document.xml.rels":
                try:
                    root = ElementTree.fromstring(data)
                    removed = 0
                    for relationship in list(root):
                        if relationship.get("Type", "").endswith("/settings") or relationship.get("Target") == "settings.xml":
                            root.remove(relationship)
                            removed += 1
                    if removed:
                        ElementTree.register_namespace("", _REL_NS)
                        data = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
                        changes.append(f"removed {removed} settings relationship(s) from {info.filename}")
                except ElementTree.ParseError:
                    pass
        repaired.append((info, data))
    return repaired, changes


def _repair_office_file(path: Path, backup_path: Path) -> dict:
    original_hash = _sha256(path)
    before = _validate_office_package(path) if path.suffix.lower() in OFFICE_EXTENSIONS else None
    result = {
        "file": path.name,
        "status": "skipped",
        "original_sha256": original_hash,
        "repaired_sha256": original_hash,
        "backup": None,
        "changes": [],
        "validation_before": before,
        "validation_after": before,
    }
    if path.suffix.lower() not in OFFICE_EXTENSIONS:
        return result
    if path.suffix.lower() != ".docx":
        result["status"] = "valid" if before and before["valid"] else "unresolved"
        return result
    try:
        with zipfile.ZipFile(path) as package:
            members, changes = _repair_docx_members(package)
    except (OSError, zipfile.BadZipFile):
        result["status"] = "unresolved"
        return result
    if not changes:
        result["status"] = "valid" if before and before["valid"] else "unresolved"
        return result

    backup_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(path, backup_path)
    with tempfile.NamedTemporaryFile(dir=path.parent, suffix=path.suffix, delete=False) as handle:
        temporary = Path(handle.name)
    try:
        with zipfile.ZipFile(temporary, "w") as target:
            for info, data in members:
                target.writestr(info, data)
        after = _validate_office_package(temporary)
        if not after["valid"]:
            result.update(status="unresolved", changes=changes, backup=str(backup_path), validation_after=after)
            return result
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)

    result.update(
        status="repaired",
        repaired_sha256=_sha256(path),
        backup=str(backup_path),
        changes=changes,
        validation_after=after,
    )
    return result


def _basenames(row: dict, key: str) -> list[str]:
    return [Path(str(value)).name for value in row.get(key) or []]


def _download_files(paths: list[str], destination: Path) -> list[str]:
    from huggingface_hub import hf_hub_download

    destination.mkdir(parents=True, exist_ok=True)
    names: list[str] = []
    for repo_path in paths:
        cached = Path(hf_hub_download(REPO_ID, str(repo_path), repo_type="dataset"))
        target = destination / Path(str(repo_path)).name
        shutil.copy2(cached, target)
        names.append(target.name)
    return names


def _repair_staged_files(paths: list[Path], *, task_id: str, role: str, backup_root: Path) -> list[dict]:
    records: list[dict] = []
    for path in paths:
        if path.suffix.lower() not in OFFICE_EXTENSIONS:
            continue
        record = _repair_office_file(path, backup_root / task_id / role / path.name)
        record.update(task_id=task_id, role=role, staged_path=str(path))
        records.append(record)
    return records


def _write_instruction(task_dir: Path, row: dict, input_names: list[str]) -> None:
    lines = [str(row.get("prompt") or "").strip(), ""]
    if input_names:
        lines.extend(["## Input files", *[f"- /home/user/{name}" for name in input_names], ""])
    expected = _basenames(row, "deliverable_files")
    if expected:
        lines.extend(["## Expected deliverable filename(s)", *[f"- {name}" for name in expected], ""])
    lines.extend(
        [
            "## Submission",
            "Save final deliverables in /home/user/deliverables and submit those files with the finish tool.",
            "",
        ]
    )
    (task_dir / "instruction.md").write_text("\n".join(lines), encoding="utf-8")


def _write_task_toml(task_dir: Path, row: dict, input_names: list[str], gold_names: list[str]) -> None:
    lines = [
        'schema_version = "1.1"',
        f'task_id = "{row["task_id"]}"',
        f'sector = """{_toml_escape(str(row.get("sector") or ""))}"""',
        f'occupation = """{_toml_escape(str(row.get("occupation") or ""))}"""',
        f"reference_files = {json.dumps(input_names)}",
        f"reference_deliverables = {json.dumps(gold_names)}",
        "",
        "[task]",
        f'name = "{row["task_id"]}"',
        "",
        "[environment]",
        'workdir = "/home/user"',
        "cpus = 4",
        "memory_mb = 16384",
        "storage_mb = 32768",
        "",
        "[agent]",
        'user = "root"',
        "timeout_sec = 14400",
        "",
        "[verifier]",
        'script = "tests/test.sh"',
        'user = "root"',
        "timeout_sec = 1800",
        "",
        "[verifier.env]",
        'GDPVAL_JUDGE_API_KEY = "${GDPVAL_JUDGE_API_KEY}"',
        'GDPVAL_JUDGE_MODEL = "${GDPVAL_JUDGE_MODEL}"',
        "",
    ]
    (task_dir / "task.toml").write_text("\n".join(lines), encoding="utf-8")


def _write_tests(task_dir: Path, row: dict) -> None:
    tests_dir = task_dir / "tests"
    tests_dir.mkdir(parents=True, exist_ok=True)
    metadata = {
        "task_id": row["task_id"],
        "prompt": row.get("prompt") or "",
        "rubric": row.get("rubric_pretty") or row.get("rubric_json") or "",
    }
    (tests_dir / "metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    (tests_dir / "evaluate.py").write_text(_VERIFIER_SOURCE, encoding="utf-8")
    test_sh = tests_dir / "test.sh"
    test_sh.write_text("#!/bin/bash\nset -euo pipefail\npython /tests/evaluate.py\n", encoding="utf-8")
    test_sh.chmod(0o755)


def _write_dataset_toml(out: Path, *, name: str, split: str, description: str, default_agent: str) -> None:
    (out / "dataset.toml").write_text(
        "\n".join(
            [
                "[dataset]",
                f'name = "{name}"',
                'type = "sandbox"',
                f'description = "{_toml_escape(description)}"',
                'default_sandbox = "docker"',
                f'default_agent = "{default_agent}"',
                f'split = "{split}"',
                "",
            ]
        ),
        encoding="utf-8",
    )


def build_benchmark(
    *,
    name: str = "gdpval",
    split: str = "train",
    out_dir: str | Path,
    catalog_entry: dict | None = None,
    limit: int | None = None,
    occupations: list[str] | None = None,
    default_agent: str = "stirrup",
    repair_office_files: bool = True,
    clean: bool = False,
    register: bool = True,
) -> Path:
    """Download and materialize the official GDPval split."""
    from datasets import load_dataset

    if catalog_entry:
        split = catalog_entry.get("eval_split") or split
        default_agent = catalog_entry.get("default_agent") or default_agent

    out = Path(out_dir).expanduser()
    if clean and out.exists():
        shutil.rmtree(out)
    out.mkdir(parents=True, exist_ok=True)

    rows = [dict(row) for row in load_dataset(REPO_ID, split=split)]
    if occupations:
        allowed = {value.casefold().strip() for value in occupations}
        rows = [row for row in rows if str(row.get("occupation", "")).casefold().strip() in allowed]
    if limit is not None:
        rows = rows[:limit]

    repair_records: list[dict] = []
    registry_rows: list[dict] = []
    backup_root = out / ".office-repair-originals"
    for row in rows:
        task_id = str(row.get("task_id") or "")
        if not task_id:
            continue
        task_dir = out / task_id
        if task_dir.exists():
            shutil.rmtree(task_dir)
        inputs_dir = task_dir / "environment" / "files"
        gold_dir = task_dir / "tests" / "reference"
        input_names = _download_files(list(row.get("reference_files") or []), inputs_dir)
        gold_names = _download_files(list(row.get("deliverable_files") or []), gold_dir)
        if repair_office_files:
            repair_records.extend(
                _repair_staged_files([inputs_dir / name for name in input_names], task_id=task_id, role="input", backup_root=backup_root)
            )
            repair_records.extend(
                _repair_staged_files([gold_dir / name for name in gold_names], task_id=task_id, role="expert", backup_root=backup_root)
            )

        (task_dir / "environment" / "Dockerfile").write_text(_DOCKERFILE, encoding="utf-8")
        _write_instruction(task_dir, row, input_names)
        _write_task_toml(task_dir, row, input_names, gold_names)
        _write_tests(task_dir, row)
        registry_rows.append(
            {
                "task_id": task_id,
                "id": task_id,
                "instruction": str(row.get("prompt") or ""),
                "task_path": str(task_dir),
                "occupation": row.get("occupation", ""),
                "sector": row.get("sector", ""),
            }
        )

    if not registry_rows:
        raise RuntimeError(f"no GDPval tasks materialized from {REPO_ID} split={split}")

    description = (catalog_entry or {}).get("description") or "GDPval: 220 economically valuable knowledge-work tasks with expert deliverables."
    _write_dataset_toml(out, name=name, split=split, description=description, default_agent=default_agent)
    if repair_office_files:
        statuses: dict[str, int] = {}
        for record in repair_records:
            statuses[record["status"]] = statuses.get(record["status"], 0) + 1
        manifest = {
            "schema_version": 1,
            "policy": "Known GDPval DOCX repairs only; unknown defects are reported as unresolved.",
            "summary": {"files": len(repair_records), "statuses": statuses},
            "files": repair_records,
        }
        (out / "office_repair_manifest.json").write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")

    if register:
        try:
            from rllm.data import DatasetRegistry

            DatasetRegistry.register_dataset(
                name=name,
                data=registry_rows,
                split=split,
                source=REPO_ID,
                description=description,
                category=(catalog_entry or {}).get("category", "agentic"),
            )
        except Exception:
            logger.warning("[gdpval] could not register rows in DatasetRegistry", exc_info=True)
    return out


_DOCKERFILE = r'''FROM python:3.10-bookworm

ENV DEBIAN_FRONTEND=noninteractive
RUN apt-get update && apt-get install -y --no-install-recommends \
    libreoffice libreoffice-writer fonts-dejavu-core tesseract-ocr tesseract-ocr-eng \
    pandoc poppler-utils ghostscript ffmpeg graphviz libgraphviz-dev openjdk-17-jre-headless \
    gdal-bin libgdal-dev libgeos-dev libproj-dev gcc g++ cmake pkg-config make gfortran \
    libsndfile1 libzbar0 libgl1-mesa-glx libglib2.0-0 libcairo2-dev libpango1.0-dev \
    libgdk-pixbuf2.0-dev libffi-dev libxml2-dev libxslt1-dev wkhtmltopdf unrar-free \
    espeak-ng texlive-latex-base texlive-latex-recommended texlive-latex-extra latexmk \
    && rm -rf /var/lib/apt/lists/*
RUN fc-cache -f
RUN pip install --no-cache-dir --upgrade 'pip>=25.0' 'setuptools<81'
ARG PIP_FLAGS="--no-cache-dir --prefer-binary --retries 10 --timeout 60"
RUN pip install $PIP_FLAGS \
    numpy 'pandas<2' scipy matplotlib 'Pillow<10' seaborn plotly bokeh numpy-financial \
    sympy h5py tables statsmodels scikit-learn scikit-image xgboost lightgbm shap \
    nltk gensim 'spacy<3.8' textblob opencv-python pytesseract qrcode pyzbar imgkit
RUN pip install $PIP_FLAGS \
    ffmpeg-python pydub 'moviepy<2' soundfile 'librosa>=0.10' mutagen \
    'python-docx<1' 'python-pptx<1' openpyxl xlrd PyMuPDF pdf2image pdfplumber \
    pypandoc docx2txt odfpy pyxlsb 'camelot-py[base]' fpdf2 'reportlab<4' weasyprint \
    graphviz 'pydot<2' networkx svglib svgwrite cairosvg wordcloud \
    shapely fiona geopandas geopy rasterio rdkit biopython
RUN pip install $PIP_FLAGS \
    markdownify anytree rarfile chardet tqdm tabulate faker loguru rapidfuzz \
    pycountry cryptography pyopenssl requests pytest
RUN pip install --no-cache-dir 'setuptools<81'
RUN mkdir -p /home/user/deliverables
WORKDIR /home/user
CMD ["tail", "-f", "/dev/null"]
'''


_VERIFIER_SOURCE = r'''from __future__ import annotations

import base64
import json
import mimetypes
import os
import random
import re
import subprocess
import tempfile
import urllib.request
from pathlib import Path

import fitz

CANDIDATE_DIR = Path("/home/user/deliverables")
REFERENCE_DIR = Path("/tests/reference")
REWARD_PATH = Path("/logs/verifier/reward.json")
MEDIA_EXTENSIONS = {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".avi", ".webm"}


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".csv", ".json", ".py", ".html", ".xml"}:
            return path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".pdf":
            with fitz.open(path) as document:
                return "\n".join(page.get_text() for page in document)
        if suffix in {".docx", ".doc"}:
            from docx import Document

            document = Document(path)
            return "\n".join([*(p.text for p in document.paragraphs), *(" | ".join(c.text for c in row.cells) for table in document.tables for row in table.rows)])
        if suffix in {".xlsx", ".xlsm"}:
            import openpyxl

            workbook = openpyxl.load_workbook(path, data_only=False, read_only=True)
            parts = []
            for sheet in workbook.worksheets:
                parts.append(f"## Sheet: {sheet.title}")
                parts.extend(" | ".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))
            return "\n".join(parts)
        if suffix == ".pptx":
            from pptx import Presentation

            deck = Presentation(path)
            return "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
        return f"[{path.name}: binary file, {path.stat().st_size} bytes]"
    except Exception as exc:
        return f"[{path.name}: extraction failed: {exc}]"


def render_images(path: Path, max_pages: int = 8) -> list[str]:
    suffix = path.suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        mime = mimetypes.guess_type(path.name)[0] or "image/png"
        return [f"data:{mime};base64,{base64.b64encode(path.read_bytes()).decode()}"]
    pdf_path = path
    temporary = None
    if suffix in {".docx", ".xlsx", ".xlsm", ".pptx"}:
        temporary = tempfile.TemporaryDirectory()
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", temporary.name, str(path)],
            check=False,
            capture_output=True,
            timeout=180,
        )
        pdf_path = Path(temporary.name) / f"{path.stem}.pdf"
    if pdf_path.suffix.lower() != ".pdf" or not pdf_path.exists():
        if temporary:
            temporary.cleanup()
        return []
    images = []
    try:
        with fitz.open(pdf_path) as document:
            for page in list(document)[:max_pages]:
                data = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).tobytes("png")
                images.append(f"data:image/png;base64,{base64.b64encode(data).decode()}")
    finally:
        if temporary:
            temporary.cleanup()
    return images


def describe(label: str, paths: list[Path]) -> tuple[str, list[str]]:
    text = [f"## {label}"]
    images = []
    for path in paths:
        text.extend([f"\n### File: {path.name}", extract_text(path)[:120000]])
        images.extend(render_images(path))
    return "\n".join(text), images[:16]


def judge(metadata: dict, candidate: list[Path], reference: list[Path]) -> tuple[float, dict]:
    candidate_text, candidate_images = describe("Submission A", candidate)
    reference_text, reference_images = describe("Submission B", reference)
    seed = int.from_bytes(metadata["task_id"].encode(), "little")
    swapped = bool(random.Random(seed).getrandbits(1))
    if swapped:
        candidate_text, reference_text = reference_text, candidate_text
        candidate_images, reference_images = reference_images, candidate_images

    content = [
        {
            "type": "text",
            "text": (
                "You are an impartial evaluator. Compare two submissions for the task and rubric below. "
                "Judge correctness, completeness, usability, and visual quality. Return only JSON: "
                '{"winner":"A"}, {"winner":"B"}, or {"winner":"tie"}.\n\n'
                f"TASK:\n{metadata['prompt']}\n\nRUBRIC:\n{metadata['rubric']}\n\n{candidate_text}"
            ),
        }
    ]
    content.extend({"type": "image_url", "image_url": {"url": image}} for image in candidate_images)
    content.append({"type": "text", "text": reference_text})
    content.extend({"type": "image_url", "image_url": {"url": image}} for image in reference_images)
    payload = {
        "model": os.environ["GDPVAL_JUDGE_MODEL"],
        "messages": [{"role": "user", "content": content}],
        "temperature": 0,
        "max_tokens": 2048,
    }
    request = urllib.request.Request(
        os.environ.get("GDPVAL_JUDGE_BASE_URL", "https://openrouter.ai/api/v1").rstrip("/") + "/chat/completions",
        data=json.dumps(payload).encode(),
        headers={"Authorization": f"Bearer {os.environ['GDPVAL_JUDGE_API_KEY']}", "Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=600) as response:
        body = json.loads(response.read())
    answer = body["choices"][0]["message"]["content"]
    match = re.search(r'"winner"\s*:\s*"(A|B|tie)"', answer, re.IGNORECASE)
    if not match:
        raise ValueError(f"judge returned no winner: {answer[:500]}")
    winner = match.group(1).lower()
    if swapped:
        score = {"a": 0.0, "b": 1.0, "tie": 0.5}[winner]
    else:
        score = {"a": 1.0, "b": 0.0, "tie": 0.5}[winner]
    return score, {"winner": winner, "positions_swapped": swapped, "judge_model": payload["model"]}


def write_reward(reward: float, *, ungraded: bool = False, metadata: dict | None = None) -> None:
    REWARD_PATH.parent.mkdir(parents=True, exist_ok=True)
    REWARD_PATH.write_text(
        json.dumps(
            {
                "reward": reward,
                "is_correct": reward > 0.5,
                "signals": {"pairwise_win": reward, "ungraded": 1.0 if ungraded else 0.0},
                "metadata": metadata or {},
            }
        )
    )


def main() -> None:
    metadata = json.loads(Path("/tests/metadata.json").read_text())
    candidate = sorted(path for path in CANDIDATE_DIR.rglob("*") if path.is_file()) if CANDIDATE_DIR.exists() else []
    reference = sorted(path for path in REFERENCE_DIR.rglob("*") if path.is_file()) if REFERENCE_DIR.exists() else []
    if not candidate:
        write_reward(0.0, metadata={"reason": "no_deliverable"})
        return
    if not reference:
        write_reward(0.0, ungraded=True, metadata={"reason": "no_expert_reference"})
        return
    if all(path.suffix.lower() in MEDIA_EXTENSIONS for path in [*candidate, *reference]):
        write_reward(0.0, ungraded=True, metadata={"reason": "media_requires_files_api"})
        return
    try:
        score, decision = judge(metadata, candidate, reference)
    except Exception as exc:
        write_reward(0.0, ungraded=True, metadata={"reason": "judge_failed", "error": str(exc)})
        return
    write_reward(score, metadata=decision)


if __name__ == "__main__":
    main()
'''
