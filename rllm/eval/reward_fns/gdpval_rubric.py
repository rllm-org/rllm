"""GDPval grader: faithful weighted-rubric LLM-as-judge.

GDPval (``openai/gdpval``) is OpenAI's benchmark of real, economically valuable
knowledge work: each task is a role-first ``prompt`` plus ``reference_files``,
and the expected output is a document *deliverable* (xlsx / docx / pptx / pdf).
Each task ships a machine-readable rubric (``rubric_json``): a list of weighted
pass/fail criteria::

    [{"score": 2, "criterion": "The workbook contains a worksheet named ...", ...},
     {"score": 5, "criterion": "The recommended sample size is 220 ...", ...},
     {"score": -2, "criterion": "The deliverable fabricates figures not in ...", ...}]

This module reproduces GDPval's *automated* grading (the reproducible LLM-judge
methodology, not the headline human pairwise win-rate): a multimodal judge
reviews the deliverable (structured text + page images) and returns a Yes/No per
rubric item — with a capped number of images in one call, re-checked in a
verification pass — and the reward is the weight-normalized sum of satisfied
criteria (giant deliverables are output-chunked and OR-aggregated).

Score aggregation (matches GDPval's rubric semantics):
    earned         = sum(c.score for c in rubric if judge_says_condition_holds(c))
    total_possible = sum(c.score for c in rubric if c.score > 0)
    reward         = clip(earned / total_possible, 0.0, 1.0)

Positive-weight criteria *add* when satisfied; negative-weight (penalty)
criteria *subtract* when their (undesirable) condition is present. Reward is
clipped to [0, 1] so penalties cannot drive it below zero.

Entry point: :func:`gdpval_rubric_reward_fn` — a host-side ``reward_fn``, the
same shape ``claw_eval`` uses, registered on the ``gdpval`` dataset. It
runs on the submission files the Stirrup harness already copied out of the
sandbox, so the rubric (which states the answers) never enters the solver's
environment, and a run can be re-graded without re-running the solver.

This is *not* AA's GDPval-AA methodology. AA grades GDPval-AA v2 by pairwise
comparison against human expert deliverables, fit with Bradley-Terry and
anchored to human experts at 1000 — no rubric involved. This module is the
absolute, dense [0, 1] signal that OpenAI's own guidance offers as a rough
estimate; the pairwise arena is a separate stage. This is the sole public
``gdpval`` evaluation mode in rLLM.

Audio/video deliverables (2 of the 220 tasks) are routed to Gemini's Files API,
which watches them natively — see :mod:`rllm.eval.reward_fns._gdpval_media`. That
path needs a direct ``GEMINI_API_KEY``; without one those tasks are reported
ungraded rather than scored from a filename.

Judge model resolution (first hit wins), mirroring ``claw_eval``:
  1. ``GDPVAL_JUDGE_MODEL`` / ``GDPVAL_JUDGE_BASE_URL`` / ``GDPVAL_JUDGE_API_KEY``.
  2. ``task.metadata`` keys ``judge_model`` / ``judge_base_url`` / ``judge_api_key``.
  3. The user's rLLM provider config (``~/.rllm/config.json``) via litellm.
"""

from __future__ import annotations

import base64
import json
import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path

from rllm.env import env_int
from rllm.eval.reward_fns._gdpval_media import (
    MediaUploadError,
    ffprobe_metadata,
    is_media,
    make_media_judge,
    media_api_key,
    media_judge_model,
    upload_media,
)
from rllm.eval.reward_fns._gdpval_media import (
    render_docker_image as _render_docker_image,
)
from rllm.eval.types import EvalOutput, Signal
from rllm.types import Episode, Task

logger = logging.getLogger(__name__)

# Deliberately no module-level ``SYSTEM_PROMPT``: this benchmark's solver prompt
# is AA's, byte-exact, and a verifier hint injected by
# ``_resolver.get_verifier_system_prompt`` would corrupt that.


# Multimodal rubric grader: the judge sees the deliverable as
# structured text + page images and returns a Yes/No per rubric item. Each
# rubric chunk gets one call (<=RUBRIC_CHUNK_SIZE items); a giant
# deliverable is split into output-chunks and OR-aggregated; a verification pass
# re-checks each Yes/No against its justification.
GRADING_SYSTEM_PROMPT = """You are a precise, fair, and strict evaluator for GDPval, a benchmark of professional knowledge work.

Your task: evaluate a model's output against a detailed rubric. The model was given a task prompt and produced output files and/or a text response. You will receive:
1. The original task prompt
2. The input files supplied with the task
3. The model's output: filenames, structured text, and attached visual/media content
4. A list of rubric items, each with a score and a criterion

For EACH rubric item, decide if the criterion is met (Yes) or not met (No).

Evaluation rules:
- Base judgments ONLY on evidence present in the model's output. Do not guess or infer.
- Each rubric item is a binary Yes/No:
  - Positive-score items (score > 0): "Yes" = the criterion IS satisfied; "No" = it is not.
  - Negative-score items (score < 0): "Yes" = the (undesirable) condition IS present; "No" = it is not.
- For file-format rubrics (e.g. "the deliverable is an Excel workbook"): check the structured-text metadata (filename, extension, sheet names).
- For content rubrics: check BOTH the text extraction AND the visual renderings.
- For visual/layout rubrics: rely on the rendered images.

Output ONLY a JSON array, one object per rubric item, in the same order given:
[{"rubric_id": 1, "status": "Yes", "justification": "brief"}, {"rubric_id": 2, "status": "No", "justification": "brief"}]"""

# The verification pass re-checks each Yes/No against its justification text only.
VERIFY_SYSTEM_PROMPT = """You are auditing a prior evaluator's rubric judgments for GDPval.

For each rubric item you are given its criterion, the evaluator's Yes/No status, and their
justification. You do NOT have the original output — decide only whether each Yes/No is
*consistent with its own justification*, and flip it if the justification clearly contradicts it.

Output ONLY a JSON array in the same order: [{"rubric_id": 1, "status": "Yes"}, ...]"""

_MAX_DELIVERABLE_CHARS = 60000
RUBRIC_CHUNK_SIZE = 100  # rubric items per judge call (most tasks < 100 -> 1 call)

#: Stand-in left in place of a deliverable that could not be parsed (a missing
#: ``[gdpval]`` extra, a corrupt file). :func:`_extraction_failed` looks for it
#: so the grader can refuse the task: handed only this, a judge answers "No" to
#: every content criterion and the run scores a plausible near-zero that looks
#: like a bad deliverable rather than a broken grader.
_EXTRACTION_ERROR_MARKER = "[extraction error"


def _extraction_failed(text: str) -> bool:
    """True when every extracted file yielded only an error marker.

    Partial success is still gradable — one unreadable attachment among several
    readable ones leaves real content to judge.
    """
    body = [line for line in text.splitlines() if line.strip() and not line.startswith("### File:")]
    return bool(body) and all(line.lstrip().startswith(_EXTRACTION_ERROR_MARKER) for line in body)


def _rubric_output_char_budget() -> int:
    """Deliverable-text budget per output-chunk; larger deliverables are split
    and OR-aggregated. Override with GDPVAL_RUBRIC_OUTPUT_CHARS."""
    return max(50_000, env_int("GDPVAL_RUBRIC_OUTPUT_CHARS", 512_000))


def _rubric_input_char_budget() -> int:
    """Total budget for the task's input attachments in one judge message; over
    budget the text is fair-shared and truncated, never dropped (see
    :func:`_find_input_view`). Override with GDPVAL_RUBRIC_INPUT_CHARS.

    Chars are a rough proxy for what fits: prose runs ~5.6 chars/token but the
    dataset's big spreadsheets run 1.7-2.2, so the same budget can mean 3x the
    tokens. Size it against :func:`_max_judge_images` — page images cost ~1090
    tokens each and are usually the larger half of the request.
    """
    return max(50_000, env_int("GDPVAL_RUBRIC_INPUT_CHARS", 512_000))


def _rubric_max_tokens() -> int:
    """Judge completion budget. A full 100-item rubric's JSON array needs far
    more than 8K, or it truncates mid-array. Override with GDPVAL_RUBRIC_MAX_TOKENS."""
    return max(1024, env_int("GDPVAL_RUBRIC_MAX_TOKENS", 65536))


def _rubric_retry_backoff() -> list[int]:
    """Per-chunk retry backoff (seconds). A malformed / all-NA response is
    retried with a reparse nudge. Override with GDPVAL_RUBRIC_RETRY_BACKOFF
    (comma-separated); empty string disables retries."""
    raw = os.environ.get("GDPVAL_RUBRIC_RETRY_BACKOFF")
    if raw is None:
        return [5, 15, 45]
    return [int(x) for x in raw.split(",") if x.strip().isdigit()]


@dataclass
class Criterion:
    score: float
    text: str
    rubric_item_id: str = ""


@dataclass
class GradeResult:
    reward: float
    earned: float
    total_possible: float
    per_criterion: list[dict] = field(default_factory=list)
    ungraded: bool = False
    reason: str = ""


# --------------------------------------------------------------------------- #
# Rubric parsing + aggregation (pure; unit-testable without a judge)
# --------------------------------------------------------------------------- #


def parse_rubric(rubric_json: str | list) -> list[Criterion]:
    """Parse ``rubric_json`` (str or already-decoded list) into Criteria."""
    items = json.loads(rubric_json) if isinstance(rubric_json, str) else rubric_json
    out: list[Criterion] = []
    for it in items or []:
        try:
            score = float(it["score"])
            text = str(it["criterion"]).strip()
        except (KeyError, TypeError, ValueError):
            continue
        if not text:
            continue
        out.append(Criterion(score=score, text=text, rubric_item_id=str(it.get("rubric_item_id", ""))))
    return out


def aggregate(criteria: list[Criterion], met: list[bool]) -> GradeResult:
    """Weight-normalized aggregation. ``met[i]`` = judge says criterion i holds."""
    total_possible = sum(c.score for c in criteria if c.score > 0)
    earned = sum(c.score for c, m in zip(criteria, met, strict=False) if m)
    if total_possible <= 0:
        return GradeResult(reward=0.0, earned=earned, total_possible=total_possible, ungraded=True, reason="no_positive_weight_criteria")
    reward = max(0.0, min(1.0, earned / total_possible))
    per = [{"score": c.score, "met": bool(m), "criterion": c.text, "rubric_item_id": c.rubric_item_id} for c, m in zip(criteria, met, strict=False)]
    return GradeResult(reward=reward, earned=earned, total_possible=total_possible, per_criterion=per)


def grade_rubric(criteria: list[Criterion], deliverable_text: str, prompt: str, judge_fn) -> GradeResult:
    """Judge every criterion via ``judge_fn(criterion_text, score, deliverable, prompt) -> bool``.

    ``judge_fn`` is injected so the aggregation is testable with a stub and the
    same logic can drive either litellm (host) or the openai client (sandbox).
    A criterion whose judge call fails is treated as NOT met (conservative).
    """
    met: list[bool] = []
    for c in criteria:
        try:
            met.append(bool(judge_fn(c.text, c.score, deliverable_text, prompt)))
        except Exception as e:  # noqa: BLE001
            logger.warning("[gdpval] judge failed for criterion %s: %s", c.rubric_item_id or c.text[:40], e)
            met.append(False)
    return aggregate(criteria, met)


# --------------------------------------------------------------------------- #
# Deliverable text extraction
# --------------------------------------------------------------------------- #


def extract_deliverable_text(path: str | Path, max_chars: int | None = _MAX_DELIVERABLE_CHARS) -> str:
    """Best-effort extraction of a deliverable file's textual content.

    Supports xlsx/xlsm, docx, pptx, pdf, zip, npz, ipynb, and text-like files.
    Optional deps (openpyxl / python-docx / python-pptx / pdfplumber / numpy)
    are imported lazily; a missing dep degrades to an empty/annotated string
    rather than raising. ``max_chars=None`` disables truncation — the rubric
    grader relies on this and lets its output-chunker split giant text instead.
    """
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ""
    ext = p.suffix.lower()
    try:
        if ext in (".xlsx", ".xlsm", ".xls"):
            text = _extract_xlsx(p)
        elif ext == ".docx":
            text = _extract_docx(p)
        elif ext == ".pptx":
            text = _extract_pptx(p)
        elif ext == ".pdf":
            text = _extract_pdf(p)
        elif ext == ".ipynb":
            text = _extract_ipynb(p)
        elif ext == ".zip":
            text = _extract_zip(p)
        elif ext == ".npz":
            text = _extract_npz(p)
        elif ext in (".csv", ".tsv", ".txt", ".md", ".json", ".html", ".xml"):
            text = p.read_text(encoding="utf-8", errors="replace")
        else:
            # Unknown binary: try utf-8, else note the type.
            try:
                text = p.read_text(encoding="utf-8")
            except (UnicodeDecodeError, OSError):
                text = f"[unsupported deliverable type: {ext or 'no extension'}]"
    except Exception as e:  # noqa: BLE001
        logger.warning("[gdpval] extraction failed for %s: %s", p.name, e)
        text = f"{_EXTRACTION_ERROR_MARKER} for {p.name}: {e}]"
    header = f"### File: {p.name}\n"
    body = text if (max_chars is None or len(text) <= max_chars) else text[:max_chars] + "\n...[truncated]..."
    return header + body


def _workbook_text(path: Path) -> tuple[str, int]:
    """Sheet-by-sheet text of a workbook, plus the count of populated cells."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    populated = 0
    for ws in wb.worksheets:
        parts.append(f"#### Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            populated += sum(1 for c in cells if c)
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts), populated


def _has_uncached_formulas(path: Path) -> bool:
    """True when the workbook holds formulas with no stored result.

    A spreadsheet cell keeps both the formula and the value last computed for
    it. Excel and LibreOffice write that value on save; a workbook produced by a
    library does not, so asking for values yields ``None`` and the cell reads
    blank. Detecting this is what makes the recalculation pass conditional
    rather than a tax on every workbook.
    """
    from openpyxl import load_workbook

    try:
        formulas = load_workbook(path, read_only=True, data_only=False)
        values = load_workbook(path, read_only=True, data_only=True)
    except Exception:  # noqa: BLE001 — malformed workbook; the caller reports it
        return False
    try:
        for sheet in formulas.worksheets:
            value_sheet = values[sheet.title]
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("=") and value_sheet[cell.coordinate].value is None:
                        return True
        return False
    finally:
        formulas.close()
        values.close()


def _extract_xlsx(p: Path) -> str:
    """Workbook text, recalculating first when formulas have no stored result.

    Without this a correct workbook can read as blank: the solver writes
    ``=SUM(B2:B10)``, nothing ever computes it, and every numeric criterion is
    judged against an empty cell. LibreOffice computes on load, so a round-trip
    through it restores the values. Reproduced against a formula-only workbook.
    """
    text, populated = _workbook_text(p)
    if not _has_uncached_formulas(p):
        return text

    recalculated = _soffice_convert(p, "xlsx")
    if not recalculated:
        logger.warning("[gdpval] %s has uncached formulas and no converter is available; cells may read blank", p.name)
        return text

    import tempfile

    with tempfile.TemporaryDirectory() as td:
        fresh = Path(td) / p.name
        fresh.write_bytes(recalculated)
        try:
            fresh_text, fresh_populated = _workbook_text(fresh)
        except Exception as e:  # noqa: BLE001 — keep the original read on any failure
            logger.warning("[gdpval] recalculated %s could not be read (%s); using the original", p.name, e)
            return text
    # Only prefer the recalculated read if it actually recovered content; a
    # conversion that lost sheets or formatting must not silently win.
    if fresh_populated > populated:
        logger.info("[gdpval] recalculated %s: %d -> %d populated cells", p.name, populated, fresh_populated)
        return fresh_text
    return text


def _extract_docx(p: Path) -> str:
    import docx  # python-docx

    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(x for x in parts if x is not None)


def _extract_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"#### Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_pdf(p: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(p)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts.append(f"#### Page {i}")
            parts.append(page.extract_text() or "")
    return "\n".join(parts)


# Zip members worth extracting recursively (reuses the same per-type extractors).
_ZIP_PARSEABLE_EXTS = {
    ".xlsx",
    ".xls",
    ".xlsm",
    ".docx",
    ".pptx",
    ".pdf",
    ".ipynb",
    ".npz",
    ".csv",
    ".tsv",
    ".txt",
    ".md",
    ".json",
    ".html",
    ".xml",
    ".py",
    ".yaml",
    ".yml",
}


def _extract_ipynb(p: Path) -> str:
    """Notebook: source of every cell plus its text/plain outputs."""
    nb = json.loads(p.read_text(encoding="utf-8", errors="replace"))
    parts: list[str] = []
    for i, cell in enumerate(nb.get("cells", []) or [], 1):
        src = cell.get("source", "")
        if isinstance(src, list):
            src = "".join(src)
        parts.append(f"#### Cell {i} [{cell.get('cell_type', '')}]\n{src}")
        for out in cell.get("outputs", []) or []:
            txt = out.get("text")
            if isinstance(txt, list):
                txt = "".join(txt)
            if txt:
                parts.append(f"[stdout]\n{txt}")
            data = out.get("data", {})
            if isinstance(data, dict) and isinstance(data.get("text/plain"), list | str):
                tp = data["text/plain"]
                parts.append("[result]\n" + ("".join(tp) if isinstance(tp, list) else tp))
    return "\n".join(parts)


def _extract_zip(p: Path) -> str:
    """List archive contents and recursively extract text from parseable members."""
    import tempfile
    import zipfile

    parts: list[str] = []
    with zipfile.ZipFile(p) as zf:
        names = [n for n in zf.namelist() if not n.endswith("/")]
        parts.append("#### Archive contents\n" + "\n".join(names))
        with tempfile.TemporaryDirectory() as td:
            for name in names:
                if Path(name).suffix.lower() not in _ZIP_PARSEABLE_EXTS:
                    continue
                try:
                    dest = Path(td) / Path(name).name
                    with zf.open(name) as fh:
                        dest.write_bytes(fh.read())
                    sub = extract_deliverable_text(dest, max_chars=None)
                    if sub:
                        nested = "\n".join(sub.splitlines()[1:])
                        nested = re.sub(r"^(#{4,}) ", lambda m: "#" + m.group(0), nested, flags=re.MULTILINE)
                        parts.append(f"#### Archive member: {name}\n{nested}")
                except Exception:  # noqa: BLE001 — skip an unreadable member
                    continue
    return "\n".join(parts)


def _extract_npz(p: Path) -> str:
    """NumPy archive: array names with shape + dtype (values not dumped)."""
    import numpy as np

    parts: list[str] = []
    with np.load(p, allow_pickle=False) as z:
        for k in z.files:
            arr = z[k]
            parts.append(f"{k}: shape={getattr(arr, 'shape', None)} dtype={getattr(arr, 'dtype', None)}")
    return "#### NPZ arrays\n" + "\n".join(parts)


# --------------------------------------------------------------------------- #
# Multimodal rendering — deliverable files -> page images for the pairwise judge
# --------------------------------------------------------------------------- #

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
_OFFICE_EXTS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".odt", ".odp", ".ods", ".rtf"}
# Raw-image downscale target: bound token cost + stay under API body limits.
_MAX_IMAGE_LONG_EDGE = 1568


@dataclass
class DeliverableView:
    """A deliverable as the judge sees it: extracted text + optional page images.

    ``images`` are base64 ``data:`` URIs (PNG/JPEG) suitable for an OpenAI-format
    ``image_url`` content part. Empty when multimodal rendering is disabled,
    unsupported for the file type, or the render toolchain is unavailable.
    """

    text: str
    images: list[str] = field(default_factory=list)
    source: str = "none"
    media: list[str] = field(default_factory=list)  # audio/video file paths (graded natively via Gemini)

    @property
    def present(self) -> bool:
        return bool(self.text or self.images or self.media)


def _multimodal_enabled() -> bool:
    """Multimodal judging is on unless GDPVAL_JUDGE_MULTIMODAL=0."""
    return os.environ.get("GDPVAL_JUDGE_MULTIMODAL", "1") == "1"


def _render_dpi() -> int:
    return max(50, env_int("GDPVAL_JUDGE_RENDER_DPI", 150))


def _max_judge_images() -> int:
    """Maximum page images sent in one Gemini grading request.

    Applied per view, so inputs and the deliverable each get up to this many.
    At ~1090 tokens per page they are usually the larger half of the request:
    500 input renders alone are 546k of Gemini's 1M window, which is what makes
    a big workbook overflow. Deliverables in the 220-task set top out at 196
    pages, so trimming this bounds runaway *input* renders — where page 400 of a
    spreadsheet grid adds nothing the extracted cell text has not already said.
    """
    return max(0, env_int("GDPVAL_JUDGE_MAX_IMAGES", 500))


def _encode_raw_image(raw: bytes, ext: str) -> str:
    """Data URI for a raw image, downscaled + JPEG-transcoded to bound token cost.

    Falls back to embedding the original bytes if Pillow is unavailable or the
    image can't be decoded.
    """
    try:
        import io

        from PIL import Image

        img = Image.open(io.BytesIO(raw))
        img.load()
        long_edge = max(img.size)
        if long_edge > _MAX_IMAGE_LONG_EDGE:
            scale = _MAX_IMAGE_LONG_EDGE / long_edge
            img = img.resize((max(1, round(img.width * scale)), max(1, round(img.height * scale))))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:  # noqa: BLE001 — best-effort; embed original on any failure
        mime = "jpeg" if ext in (".jpg", ".jpeg") else (ext.lstrip(".") or "png")
        return f"data:image/{mime};base64," + base64.b64encode(raw).decode("ascii")


def _soffice_convert(p: Path, target_ext: str) -> bytes | None:
    """Convert with LibreOffice from the pinned GDPval Docker image.

    The input directory is mounted read-only and a temporary output directory
    is mounted writable. Pinning the complete image (rather than discovering a
    host LibreOffice binary) keeps LibreOffice, fonts, and system libraries the
    same on every evaluator. Returns ``None`` when Docker or conversion fails.
    """
    import subprocess
    import tempfile

    with tempfile.TemporaryDirectory() as td:
        out = Path(td) / f"{p.stem}.{target_ext}"
        src_file = p.resolve()
        try:
            from rllm.data.gdpval_aa import AA_PLATFORM

            result = subprocess.run(
                [
                    "docker",
                    "run",
                    "--rm",
                    "--platform",
                    AA_PLATFORM,
                    "--network",
                    "none",
                    "--mount",
                    f"type=bind,source={src_file.parent},target=/input,readonly",
                    "--mount",
                    f"type=bind,source={td},target=/output",
                    _render_docker_image(),
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    target_ext,
                    "--outdir",
                    "/output",
                    f"/input/{src_file.name}",
                ],
                check=True,
                capture_output=True,
                timeout=180,
            )
            if out.exists():
                return out.read_bytes()
            # LibreOffice exits 0 even when it refuses the document ("Error:
            # source file could not be loaded" on stdout), so ``check=True``
            # never fires and the missing output file is the only evidence.
            # Several upstream GDPval DOCX are damaged this way: python-docx
            # still reads their text, so the judge gets prose but no page
            # images — a silent downgrade unless it is said out loud here.
            # The reason lands on stderr ("Error: source file could not be
            # loaded"); stdout carries the javaldx noise. Read both.
            streams = b"\n".join(s for s in (result.stderr, result.stdout) if s)
            detail = [ln for ln in streams.decode("utf-8", errors="replace").splitlines() if ln.strip() and "javaldx" not in ln]
            logger.warning(
                "[gdpval] LibreOffice produced no %s for %s (exit 0, no output file): %s",
                target_ext,
                p.name,
                detail[-1] if detail else "no diagnostic output",
            )
        except (subprocess.SubprocessError, OSError) as e:
            logger.warning("[gdpval] Docker LibreOffice conversion failed for %s -> %s: %s", p.name, target_ext, e)
        return None


def _office_to_pdf(p: Path) -> bytes | None:
    """Office document as PDF bytes, or None when no converter is available.

    None means the deliverable is judged on text alone.
    """
    return _soffice_convert(p, "pdf")


def _render_pdf_pages(pdf_bytes: bytes, max_pages: int) -> list[str]:
    """Render the first ``max_pages`` PDF pages to PNG data URIs via PyMuPDF."""
    import fitz  # PyMuPDF

    uris: list[str] = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pix = page.get_pixmap(dpi=_render_dpi())
            uris.append("data:image/png;base64," + base64.b64encode(pix.tobytes("png")).decode("ascii"))
    finally:
        doc.close()
    return uris


def render_deliverable_images(path: str | Path, max_pages: int | None = None) -> list[str]:
    """Render a deliverable file to page-image data URIs for the multimodal judge.

    PDFs and office docs (via LibreOffice) are rasterized page-by-page; raster
    images are embedded directly. Returns [] — never raises — when disabled,
    unsupported, or the render toolchain (LibreOffice / PyMuPDF) is missing, so
    the judge degrades to text-only.
    """
    if not _multimodal_enabled():
        return []
    p = Path(path)
    if not p.exists() or not p.is_file():
        return []
    pages = max_pages or _max_judge_images()
    ext = p.suffix.lower()
    try:
        if ext in _IMAGE_EXTS:
            return [_encode_raw_image(p.read_bytes(), ext)]
        pdf_bytes: bytes | None = None
        if ext == ".pdf":
            pdf_bytes = p.read_bytes()
        elif ext in _OFFICE_EXTS:
            pdf_bytes = _office_to_pdf(p)
        if pdf_bytes:
            return _render_pdf_pages(pdf_bytes, pages)
    except Exception as e:  # noqa: BLE001 — rendering is best-effort; never break grading
        logger.warning("[gdpval] image render failed for %s: %s", p.name, e)
    return []


# --------------------------------------------------------------------------- #
# Host-side reward_fn entry point
# --------------------------------------------------------------------------- #


def _find_deliverable_text(task: Task, episode: Episode) -> tuple[str, str]:
    """Locate the deliverable's text. Returns (text, source_description).

    FAIL-CLOSED: the grader only ever reads a deliverable the harness has
    explicitly surfaced — never the agent's transcript or final answer. This
    keeps the judge from grading the *process* instead of the *artifact* (and
    from being gamed by an agent that describes a good deliverable it never
    produced). If nothing is surfaced, returns ("", "none") and the caller
    marks the task ungraded rather than scoring it.

    Resolution order:
      1. ``episode.artifacts['deliverable_text']`` — deliverable text surfaced directly.
      2. ``episode.artifacts['deliverable_path']`` / ``output_files`` — a file
         path (sandbox workdir mount or host) to extract.
    """
    arts = episode.artifacts or {}
    if arts.get("deliverable_text"):
        return str(arts["deliverable_text"]), "artifacts.deliverable_text"

    for c in _rank_candidates(_deliverable_candidate_paths(episode), task):
        text = extract_deliverable_text(c)
        if text:
            return text, f"file:{Path(c).name}"

    return "", "none"


def _deliverable_candidate_paths(episode: Episode) -> list[str]:
    """Surfaced deliverable file paths, in resolution order, de-duplicated.

    ``deliverables`` is what :mod:`rllm.harnesses.stirrup` publishes: the local
    paths of the files the solver named in its ``finish`` call, already copied
    out of the sandbox. ``deliverable_path`` / ``output_files`` are the older
    cookbook keys, kept so an episode from either harness grades the same.
    """
    arts = episode.artifacts or {}
    raw: list[str] = []
    for key in ("deliverables", "deliverable_path", "output_files"):
        value = arts.get(key)
        if isinstance(value, list | tuple):
            raw.extend(str(x) for x in value)
        elif value:
            raw.append(str(value))
    seen: set[str] = set()
    return [p for p in raw if not (p in seen or seen.add(p))]


def _rank_candidates(paths: list[str], task: Task) -> list[str]:
    """Order candidate deliverables so the one matching the task's *expected*
    deliverable comes first — the model may emit several files (e.g. an updated
    input alongside the real deliverable), but pairwise must compare like-for-like
    against the single gold. Prefers an expected-basename match, then a matching
    file extension; ties keep surfaced order (stable sort).
    """
    meta = task.metadata or {}
    expected = meta.get("expected_deliverables") or meta.get("reference_deliverables") or []
    exp_stems = [Path(str(n)).stem.lower() for n in expected]
    exp_exts = {Path(str(n)).suffix.lower() for n in expected if Path(str(n)).suffix}

    def rank(p: str) -> int:
        pth = Path(p)
        stem = pth.stem.lower()
        score = 0
        if any(e and (e in stem or stem in e) for e in exp_stems):
            score -= 2  # basename matches the expected deliverable
        if pth.suffix.lower() in exp_exts:
            score -= 1  # right file type
        return score

    return sorted(paths, key=rank)


def _find_deliverable_view(task: Task, episode: Episode) -> DeliverableView:
    """Deliverable as (combined text + page images), fail-closed.

    Grades ALL surfaced deliverable files (not just the top-ranked one) — the
    model may emit several files that together satisfy the rubric. Text is
    concatenated (untruncated; the grader's output-chunker splits it).

    Images preserve file order and page order. With Gemini as the judge, the
    high default cap lets normal GDPval submissions fit in one request without
    round-robin allocation or image batching.
    """
    arts = episode.artifacts or {}
    if arts.get("deliverable_text"):
        return DeliverableView(str(arts["deliverable_text"]), [], "artifacts.deliverable_text")

    texts: list[str] = []
    names: list[str] = []
    media: list[str] = []
    per_file: list[list[str]] = []
    for c in _rank_candidates(_deliverable_candidate_paths(episode), task):
        if is_media(c):
            if Path(c).is_file():
                media.append(c)
                names.append(Path(c).name)
                # ffprobe costs a docker spin; skip it when the task will be
                # ungraded anyway for want of a key.
                if media_api_key():
                    texts.append(ffprobe_metadata(c))
            continue
        text = extract_deliverable_text(c, max_chars=None)
        if not text:
            continue
        rendered = render_deliverable_images(c)
        if rendered:
            text += f"\n\n[Attached visual content: {Path(c).name} is rendered as {len(rendered)} consecutive page image(s) below.]"
        texts.append(text)
        names.append(Path(c).name)
        per_file.append(rendered)

    images = _flatten_images(per_file, _max_judge_images())
    if not texts and not images and not media:
        return DeliverableView("", [], "none")
    return DeliverableView("\n\n".join(texts), images, "file:" + ",".join(names), media)


def _flatten_images(per_file: list[list[str]], cap: int) -> list[str]:
    """Flatten page images in file/page order and apply the single-call cap."""
    total = sum(len(pages) for pages in per_file)
    picked = [image for pages in per_file for image in pages][:cap]
    if total > len(picked):
        logger.warning(
            "[gdpval] %d page images available across %d file(s); carrying %d (GDPVAL_JUDGE_MAX_IMAGES)",
            total,
            len(per_file),
            len(picked),
        )
    return picked


def _input_file_paths(task: Task) -> list[str]:
    """On-disk paths of the task's input attachments (``reference_files`` staged
    under ``<task_dir>/environment/files/`` by the builder).

    ``reference_files`` holds *in-sandbox absolute* paths (``/home/user/x.docx``)
    because the solver prompt must quote the paths the solver will see. Grading
    runs on the host, so only the basename carries over — and it has to be taken
    explicitly: ``Path(base) / "environment" / "files" / "/home/user/x.docx"``
    discards ``base`` and hands back the sandbox path, which never exists on the
    host. That silently drops every input file from the judge's context, so the
    judge scores a deliverable against source material it was never shown.
    """
    meta = task.metadata or {}
    names = meta.get("reference_files") or []
    if not names or task.dataset_dir is None:
        return []
    base = task.dataset_dir / task.sub_dir if task.sub_dir else task.dataset_dir
    staged = Path(base) / "environment" / "files"
    return [str(staged / Path(n).name) for n in names]


def _fair_share(lengths: list[int], budget: int) -> list[int]:
    """Per-file character allowances totalling at most *budget*.

    Water-filling: every file is offered an equal share, and each one needing
    less than its share keeps all of it and releases the remainder to the rest.
    A 2 KB memo beside a 1 MB workbook is therefore never cut to make room —
    it fits, and the workbook absorbs what is left. Mirrors how page images are
    already handed out (round-robin per file) rather than first-come.
    """
    allow = [0] * len(lengths)
    remaining = budget
    pending = set(range(len(lengths)))
    while pending:
        share = remaining // len(pending)
        fits = {i for i in pending if lengths[i] <= share}
        if not fits:  # every survivor wants more than its share — split evenly
            for i in pending:
                allow[i] = share
            break
        for i in fits:
            allow[i] = lengths[i]
            remaining -= lengths[i]
        pending -= fits
    return allow


def _find_input_view(task: Task) -> DeliverableView:
    """Input attachments as structured text and consecutive page images.

    AA provides reference files to its judge as both text and images. As with
    model outputs, each input file's pages stay consecutive and are held only
    in memory.

    Over budget, the text is truncated to fit rather than dropped. Dropping
    made a single oversized attachment silently remove *every* input from the
    judge's context, so the deliverable got graded against source material the
    judge never saw — and the resulting score looked no different from an
    honest one. A truncated workbook still carries its headers and most of its
    rows, and the marker tells the judge the tail is missing. Page images are
    kept either way; they are capped separately by GDPVAL_JUDGE_MAX_IMAGES.
    """
    budget = _rubric_input_char_budget()
    bodies: list[str] = []
    notes: list[str] = []
    per_file: list[list[str]] = []
    names: list[str] = []
    for path in _input_file_paths(task):
        t = extract_deliverable_text(path, max_chars=None)
        if not t:
            continue
        rendered = render_deliverable_images(path)
        bodies.append(t)
        notes.append(f"\n\n[Attached visual content: {Path(path).name} is rendered as {len(rendered)} consecutive page image(s) below.]" if rendered else "")
        per_file.append(rendered)
        names.append(Path(path).name)
    if not bodies:
        return DeliverableView("", [], "none")
    total = sum(len(b) for b in bodies)
    if total > budget:
        allow = _fair_share([len(b) for b in bodies], budget)
        logger.warning(
            "[gdpval] input attachments %d chars > budget %d; truncating to fit (per-file: %s)",
            total,
            budget,
            ", ".join(f"{n}={a}" for n, a in zip(names, allow, strict=True)),
        )
        bodies = [b if len(b) <= a else b[:a] + "\n...[truncated]..." for b, a in zip(bodies, allow, strict=True)]
    sections = [b + n for b, n in zip(bodies, notes, strict=True)]
    images = _flatten_images(per_file, _max_judge_images())
    return DeliverableView("\n\n".join(sections), images, "file:" + ",".join(names))


def _resolve_judge(task: Task) -> tuple[str, str | None, str | None]:
    """Return (model, base_url, api_key); mirrors claw_eval's resolution."""
    env_model = os.environ.get("GDPVAL_JUDGE_MODEL")
    if env_model:
        return env_model, os.environ.get("GDPVAL_JUDGE_BASE_URL"), os.environ.get("GDPVAL_JUDGE_API_KEY")

    meta_model = task.metadata.get("judge_model")
    if meta_model:
        return meta_model, task.metadata.get("judge_base_url"), task.metadata.get("judge_api_key")

    try:
        from rllm.eval.config import get_provider_info, load_config

        cfg = load_config()
        if cfg.provider == "custom":
            return cfg.model, cfg.base_url or None, cfg.api_key or "EMPTY"
        info = get_provider_info(cfg.provider)
        if info and cfg.model:
            prefix = info.litellm_prefix
            model = f"{prefix}/{cfg.model}" if prefix and not cfg.model.startswith(prefix + "/") else cfg.model
            return model, None, cfg.api_key or None
    except Exception:
        logger.debug("[gdpval] could not resolve judge from rLLM config", exc_info=True)
    return "", None, None


def _make_rubric_judge(model: str, base_url: str | None, api_key: str | None):
    """Return ``call(messages) -> text`` for the multimodal rubric judge.

    drop_params lets litellm silently drop provider-unsupported params (reasoning
    models reject temperature=0.0); a multimodal (image) message that a non-vision
    model rejects falls back to the text-only messages.
    """
    import litellm

    def call(messages: list[dict], text_only_messages: list[dict] | None = None) -> str:
        kwargs: dict = {"model": model, "temperature": 0.0, "drop_params": True, "max_tokens": _rubric_max_tokens()}
        if base_url:
            kwargs["base_url"] = base_url
        if api_key:
            kwargs["api_key"] = api_key
        try:
            resp = litellm.completion(messages=messages, **kwargs)
        except Exception as e:  # noqa: BLE001 — non-vision model / image rejection
            if text_only_messages is None:
                raise
            # Record the downgrade. Without this the run still reports the image
            # count it *prepared*, so a non-vision judge silently looks like it
            # graded the page renderings — and visual criteria were decided from
            # text alone.
            call.text_only_fallbacks += 1
            logger.warning("[gdpval] multimodal rubric call failed (%s); retrying text-only", e)
            resp = litellm.completion(messages=text_only_messages, **kwargs)
        return resp.choices[0].message.content or ""

    call.text_only_fallbacks = 0
    return call


def _rubric_lines(chunk: list[tuple[int, Criterion]]) -> str:
    return "\n\n".join(f"**Rubric {i}** (score: {int(c.score):+d})\n{c.text}" for i, c in chunk)


def _grading_messages(
    prompt: str,
    text: str,
    images: list[str],
    chunk: list[tuple[int, Criterion]],
    input_text: str = "",
    input_images: list[str] | None = None,
) -> tuple[list[dict], list[dict]]:
    """Multimodal + text-only messages for grading one rubric chunk against a
    deliverable (structured text + page images), optionally alongside the task's
    input attachments."""
    task_section = f"## 1. Task Prompt\n{prompt}"
    input_images = input_images or []
    input_notes = ""
    if input_images:
        input_notes = f"\n\n- {len(input_images)} input page image(s) are attached immediately below, in file order with consecutive pages."
    input_section = f"## 2. Input Files\n{input_text or 'No input files were provided.'}{input_notes}"
    attachment_notes: list[str] = []
    if images:
        attachment_notes.append(f"{len(images)} page image(s) are attached immediately below, in file order with consecutive pages.")
    if "# Media file:" in text:
        attachment_notes.append("The audio/video file(s) named above are attached through the Gemini Files API.")
    notes = "\n\n" + "\n".join(f"- {note}" for note in attachment_notes) if attachment_notes else ""
    output_section = f"## 3. Model Output\n{text}{notes}"
    rubric_section = f"## 4. Rubric Items\n{_rubric_lines(chunk)}\n\nEvaluate each rubric item. Output ONLY the JSON array."
    body = "\n\n---\n\n".join((task_section, input_section, output_section, rubric_section))
    text_only = [{"role": "system", "content": GRADING_SYSTEM_PROMPT}, {"role": "user", "content": body}]
    separator = "\n\n---\n\n"
    parts: list[dict] = [
        {"type": "text", "text": task_section + separator},
        {"type": "text", "text": input_section + "\n\n"},
    ]
    for uri in input_images:
        parts.append({"type": "image_url", "image_url": {"url": uri}})
    parts.append({"type": "text", "text": separator + output_section + "\n\n"})
    for uri in images:
        parts.append({"type": "image_url", "image_url": {"url": uri}})
    parts.append({"type": "text", "text": separator + rubric_section})
    multimodal = [{"role": "system", "content": GRADING_SYSTEM_PROMPT}, {"role": "user", "content": parts}]
    return multimodal, text_only


def _clean_and_parse_json(json_str: str):
    """Parse JSON tolerantly: strip trailing commas, normalize smart quotes, and
    close a truncated array. Raises ValueError if nothing parses."""
    if not json_str or not json_str.strip():
        raise ValueError("empty JSON string")
    s = json_str.strip()
    for cand in (s, re.sub(r",\s*]", "]", re.sub(r",\s*}", "}", s)), s.replace("\u201c", '"').replace("\u201d", '"')):
        try:
            return json.loads(cand)
        except json.JSONDecodeError:
            continue
    if s.startswith("[") and not s.endswith("]"):
        # Truncated array (e.g. the judge hit its token limit mid-output): close
        # it, dropping any dangling trailing comma / partial final element.
        for tail in (s, s.rstrip().rstrip(","), s[: s.rfind("}") + 1] if "}" in s else s):
            try:
                return json.loads(tail + "]")
            except json.JSONDecodeError:
                continue
    raise ValueError("could not parse JSON")


def _parse_grading_response(text: str, ids: list[int]) -> dict[int, dict]:
    """Parse the judge's JSON array -> {rubric_id: {hit, justification}}; missing
    ids become NA (conservative). Tolerates fenced/malformed/truncated output."""
    out: dict[int, dict] = {i: {"hit": "NA", "justification": ""} for i in ids}
    data = None
    m = re.search(r"```(?:json)?\s*(\[[\s\S]*?\])\s*```", text, re.DOTALL)
    if m:
        try:
            data = _clean_and_parse_json(m.group(1))
        except (ValueError, json.JSONDecodeError):
            data = None
    if data is None:
        first, last = text.find("["), text.rfind("]")
        if first != -1 and last > first:
            segment = text[first : last + 1]
        elif first != -1:
            segment = text[first:]  # truncated: no closing bracket present
        else:
            segment = text
        try:
            data = _clean_and_parse_json(segment)
        except (ValueError, json.JSONDecodeError):
            data = None
    for item in data or []:
        if not isinstance(item, dict):
            continue
        try:
            rid = int(item.get("rubric_id"))
        except (TypeError, ValueError):
            continue
        if rid in out:
            status = str(item.get("status", "")).strip().lower()
            out[rid] = {"hit": 1 if status in ("yes", "y", "true", "1") else 0, "justification": str(item.get("justification", ""))}
    return out


_REPARSE_NUDGE = {
    "role": "user",
    "content": "Your previous response could not be parsed. Output ONLY a valid JSON array, no markdown fences, no extra text.",
}


def _grade_chunk_with_retries(mm: list[dict], txt: list[dict], ids: list[int], judge_call) -> dict[int, dict]:
    """One rubric chunk with retries: a malformed / all-NA response is retried
    with a reparse nudge and backoff; the last (all-NA) result is returned if
    every attempt fails (conservative -> those items count as not-met)."""
    import time

    backoff = _rubric_retry_backoff()
    last = {i: {"hit": "NA", "justification": "Parse failure"} for i in ids}
    for attempt in range(len(backoff) + 1):
        call_mm = mm if attempt == 0 else mm + [_REPARSE_NUDGE]
        call_txt = txt if attempt == 0 else txt + [_REPARSE_NUDGE]
        try:
            parsed = _parse_grading_response(judge_call(call_mm, call_txt), ids)
            if sum(1 for i in ids if parsed[i]["hit"] == "NA") < len(ids):
                return parsed  # at least one item parsed -> accept
            last = parsed
        except Exception as e:  # noqa: BLE001 — transient judge/parse error: retry
            logger.warning("[gdpval] rubric chunk attempt %d/%d failed: %s", attempt + 1, len(backoff) + 1, e)
        if attempt < len(backoff):
            time.sleep(backoff[attempt])
    return last


def _grade_one_segment(
    criteria: list[Criterion],
    text: str,
    images: list[str],
    prompt: str,
    judge_call,
    input_text: str = "",
    input_images: list[str] | None = None,
) -> dict[int, dict]:
    """Grade every criterion against one deliverable text segment.

    The rubric is chunked into <=RUBRIC_CHUNK_SIZE items per call. Every call
    receives all retained page images together.
    """
    results: dict[int, dict] = {}
    indexed = list(enumerate(criteria, 1))
    for start in range(0, len(indexed), RUBRIC_CHUNK_SIZE):
        chunk = indexed[start : start + RUBRIC_CHUNK_SIZE]
        ids = [i for i, _ in chunk]
        mm, txt = _grading_messages(prompt, text, images, chunk, input_text, input_images)
        results.update(_grade_chunk_with_retries(mm, txt, ids, judge_call))
    return results


def _grade_rubric_multimodal(
    criteria: list[Criterion],
    view: DeliverableView,
    prompt: str,
    judge_call,
    input_text: str = "",
    input_images: list[str] | None = None,
) -> tuple[list[bool], list[str]]:
    """Grade criteria against the deliverable (text + images). Splits a giant
    deliverable into output-chunks and OR-aggregates (Yes in any segment -> Yes).
    Returns (met, justifications) aligned to ``criteria``.
    """
    budget = _rubric_output_char_budget()
    text = view.text
    segments = [text[i : i + budget] for i in range(0, len(text), budget)] or [""]
    if len(segments) > 1:
        logger.info("[gdpval] deliverable text %d chars -> %d output-chunks (OR-aggregated)", len(text), len(segments))
    retained = _flatten_images([input_images or [], view.images], _max_judge_images())
    input_count = min(len(input_images or []), len(retained))
    retained_input_images = retained[:input_count]
    images = retained[input_count:]

    merged: dict[int, dict] = {}
    for seg in segments:
        seg_res = _grade_one_segment(criteria, seg, images, prompt, judge_call, input_text, retained_input_images)
        for rid, r in seg_res.items():
            prev = merged.get(rid)
            # OR-aggregate: a Yes in any segment wins; keep the Yes justification.
            if prev is None or (r["hit"] == 1 and prev["hit"] != 1):
                merged[rid] = r
    met = [merged.get(i, {"hit": "NA"})["hit"] == 1 for i in range(1, len(criteria) + 1)]
    just = [merged.get(i, {"justification": ""})["justification"] for i in range(1, len(criteria) + 1)]
    return met, just


def _verify_rubric(criteria: list[Criterion], met: list[bool], justifications: list[str], judge_call) -> list[bool]:
    """Verification pass: re-check each Yes/No against its justification and flip
    clear contradictions. Best-effort — returns the original ``met`` on failure."""
    try:
        lines = [
            f'{{"rubric_id": {i}, "criterion": {json.dumps(c.text)}, "status": "{"Yes" if m else "No"}", "justification": {json.dumps(j)}}}'
            for i, (c, m, j) in enumerate(zip(criteria, met, justifications, strict=False), 1)
        ]
        msg = [
            {"role": "system", "content": VERIFY_SYSTEM_PROMPT},
            {"role": "user", "content": "Re-check these judgments. Output ONLY the corrected JSON array.\n\n[" + ",\n".join(lines) + "]"},
        ]
        parsed = _parse_grading_response(judge_call(msg, msg), list(range(1, len(criteria) + 1)))
        return [parsed.get(i, {"hit": 1 if met[i - 1] else 0})["hit"] == 1 if parsed.get(i, {}).get("hit") != "NA" else met[i - 1] for i in range(1, len(criteria) + 1)]
    except Exception as e:  # noqa: BLE001 — verification is best-effort
        logger.warning("[gdpval] rubric verification pass failed: %s", e)
        return met


def _prepare_media_judge(media_paths: list[str], *, max_tokens: int = 65536):
    """Route a media task to the direct-Gemini judge, fail-closed.

    Returns ``(judge_call, judge_model, None)`` on success, or
    ``(None, None, ungraded_reason)`` when the task must not be graded: no
    ``GEMINI_API_KEY``, or an upload that never reached ACTIVE. Never score
    audio or video the judge could not actually watch or listen to.
    """
    key = media_api_key()
    if not key:
        logger.warning(
            "[gdpval] media deliverable but GEMINI_API_KEY unset — task will be ungraded (%s)",
            ", ".join(Path(m).name for m in media_paths),
        )
        return None, None, "media_requires_gemini_key"
    try:
        files = [upload_media(m, key) for m in media_paths]
    except MediaUploadError as e:
        logger.warning("[gdpval] media upload failed — task will be ungraded: %s", e)
        return None, None, "media_upload_failed"
    model = media_judge_model()
    return make_media_judge(model, key, files, max_tokens=max_tokens), model, None


def _structural_signals(episode: Episode) -> list[Signal]:
    """The signals the in-sandbox structural verifier would have reported.

    A catalog ``reward_fn`` *replaces* the per-task verifier (``--evaluator`` >
    catalog > per-task), so grading with the rubric would otherwise silently
    drop ``finish_called`` / ``abandoned`` / ``submission_valid`` /
    ``artifact_count``. They are recoverable host-side from the manifest the
    harness wrote, so carry them rather than lose them: a rubric score of 0 for
    "the solver abandoned" and for "the solver produced a bad deliverable" are
    different facts.
    """
    arts = episode.artifacts or {}
    manifest_path = arts.get("submission_manifest")
    manifest: dict = {}
    if manifest_path:
        try:
            manifest = json.loads(Path(str(manifest_path)).read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            logger.debug("[gdpval] no readable submission manifest at %s", manifest_path)
    termination = manifest.get("termination") or {}
    kind = termination.get("type")
    artifacts = manifest.get("artifacts") or []
    return [
        Signal(name="finish_called", value=1.0 if kind == "finish" else 0.0),
        Signal(name="abandoned", value=1.0 if kind == "abandon_task_finish" else 0.0),
        Signal(name="submission_valid", value=1.0 if (kind == "finish" and artifacts and not manifest.get("rejected_paths")) else 0.0),
        Signal(name="artifact_count", value=float(len(artifacts))),
    ]


def evaluate_rubric(task: Task, episode: Episode) -> EvalOutput:
    """Weighted-rubric LLM-judge grade over the harness-surfaced deliverable."""
    structural = _structural_signals(episode)
    rubric_raw = task.metadata.get("rubric_json") or task.metadata.get("rubric")
    criteria = parse_rubric(rubric_raw) if rubric_raw else []
    if not criteria:
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0), Signal(name="ungraded", value=1.0), *structural],
            metadata={"reason": "no_rubric", "ungraded": True},
        )

    view = _find_deliverable_view(task, episode)  # text + page images (multimodal)
    if not view.present:
        # Fail-closed: no deliverable was surfaced (never grade the transcript).
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0), Signal(name="ungraded", value=1.0), *structural],
            metadata={"reason": "no_deliverable_surfaced", "ungraded": True},
        )
    if not view.media and not view.images and _extraction_failed(view.text):
        # A deliverable exists but nothing could be read out of it — usually the
        # ``rllm[gdpval]`` extra is missing. Grading here would hand the judge an
        # error string and turn a broken host into a near-zero model score.
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0), Signal(name="ungraded", value=1.0), *structural],
            metadata={
                "reason": "deliverable_unreadable",
                "ungraded": True,
                "deliverable_source": view.source,
                "hint": "install the document parsers: uv pip install -e '.[gdpval]'",
            },
        )
    prompt = task.metadata.get("prompt") or (task.instruction if isinstance(task.instruction, str) else "") or ""

    if view.media:
        # Audio/video deliverable: grade natively via Gemini, or not at all.
        judge_call, model, ungraded_reason = _prepare_media_judge(view.media, max_tokens=_rubric_max_tokens())
        if ungraded_reason:
            return EvalOutput(
                reward=0.0,
                is_correct=False,
                signals=[Signal(name="rubric_score", value=0.0), Signal(name="ungraded_media", value=1.0), Signal(name="ungraded", value=1.0), *structural],
                metadata={"reason": ungraded_reason, "ungraded": True, "media_files": [Path(m).name for m in view.media]},
            )
    else:
        model, base_url, api_key = _resolve_judge(task)
        if not model:
            return EvalOutput(
                reward=0.0,
                is_correct=False,
                signals=[Signal(name="rubric_score", value=0.0), Signal(name="ungraded", value=1.0), *structural],
                metadata={"reason": "no_judge_configured", "ungraded": True},
            )
        judge_call = _make_rubric_judge(model, base_url, api_key)

    input_view = _find_input_view(task)
    met, justifications = _grade_rubric_multimodal(
        criteria,
        view,
        prompt,
        judge_call,
        input_view.text,
        input_view.images,
    )
    met = _verify_rubric(criteria, met, justifications, judge_call)
    result = aggregate(criteria, met)
    # What the judge actually consumed, not what was prepared for it: a
    # non-vision judge downgrades to text and every visual criterion is then
    # decided without seeing the document.
    input_images_prepared = min(len(input_view.images), _max_judge_images())
    images_prepared = min(len(view.images), max(0, _max_judge_images() - input_images_prepared))
    total_images_prepared = input_images_prepared + images_prepared
    fell_back = getattr(judge_call, "text_only_fallbacks", 0) > 0
    images_used = 0 if fell_back else images_prepared
    input_images_used = 0 if fell_back else input_images_prepared
    _write_judge_decisions(
        episode,
        {
            "task_id": str(task.id),
            "judge_model": model,
            "reward": float(result.reward),
            "earned": result.earned,
            "total_possible": result.total_possible,
            "criteria": [{"score": c.score, "criterion": c.text, "met": bool(m), "justification": j} for c, m, j in zip(criteria, met, justifications, strict=False)],
        },
    )
    return EvalOutput(
        reward=float(result.reward),
        is_correct=result.reward >= 0.5,
        signals=[Signal(name="rubric_score", value=float(result.reward)), *structural],
        metadata={
            "judge_model": model,
            "deliverable_source": view.source,
            "deliverable_images": images_used,
            "deliverable_images_prepared": images_prepared,
            "input_images": input_images_used,
            "input_images_prepared": input_images_prepared,
            "judge_images": 0 if fell_back else total_images_prepared,
            "judge_images_prepared": total_images_prepared,
            "vision_used": bool(images_used or input_images_used),
            "judge_text_only_fallback": fell_back,
            "media_files": [Path(m).name for m in view.media],
            "input_files_included": bool(input_view.text or input_view.images),
            "earned": result.earned,
            "total_possible": result.total_possible,
            "criteria_met": sum(1 for m in met if m),
            "criteria_total": len(criteria),
            "occupation": task.metadata.get("occupation", ""),
            "sector": task.metadata.get("sector", ""),
        },
    )


#: ``datasets.json``'s ``reward_fn: gdpval_rubric_reward_fn`` resolves through
#: :data:`rllm.eval.evaluator_loader._EVALUATOR_REGISTRY` to
#: :func:`evaluate_rubric` — no module-level alias needed.
evaluate = evaluate_rubric


def _write_judge_decisions(episode: Episode, payload: dict) -> None:
    """Write per-criterion decisions to ``judge_decisions.json`` in the run's
    corpus directory, beside the submitted files, so a grade can be audited
    against what the model actually produced. Best-effort.

    Prefers the harness's ``submission_dir`` (the run-keyed corpus directory)
    and falls back to the parent of a surfaced deliverable path.
    """
    arts = episode.artifacts or {}
    out_dir: Path | None = None
    if arts.get("submission_dir"):
        out_dir = Path(str(arts["submission_dir"]))
    elif arts.get("deliverable_path"):
        out_dir = Path(str(arts["deliverable_path"])).parent
    if out_dir is None:
        return
    try:
        if out_dir.is_dir():
            (out_dir / "judge_decisions.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    except OSError as e:
        logger.warning("[gdpval] could not write judge_decisions.json: %s", e)
