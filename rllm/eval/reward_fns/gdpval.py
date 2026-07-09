"""GDPval grader: faithful weighted-rubric LLM-as-judge.

GDPval (``openai/gdpval``) is OpenAI's benchmark of real, economically valuable
knowledge work: each task is a role-first ``prompt`` plus ``reference_files``,
and the expected output is a document *deliverable* (xlsx / docx / pptx / pdf).
Each task ships a machine-readable rubric (``rubric_json``): a list of weighted
pass/fail criteria::

    [{"score": 2, "criterion": "The workbook contains a worksheet named ...", ...},
     {"score": 5, "criterion": "The recommended sample size is 220 ...", ...},
     {"score": -2, "criterion": "The deliverable fabricates figures not in ...", ...}]

This module reproduces GDPval's *automated* grading (the reproducible LLM-judge
methodology, not the headline human pairwise win-rate): an LLM judge decides,
per criterion, whether the deliverable satisfies it, and the reward is the
weight-normalized sum of satisfied criteria.

Score aggregation (matches GDPval's rubric semantics):
    earned         = sum(c.score for c in rubric if judge_says_condition_holds(c))
    total_possible = sum(c.score for c in rubric if c.score > 0)
    reward         = clip(earned / total_possible, 0.0, 1.0)

Positive-weight criteria *add* when satisfied; negative-weight (penalty)
criteria *subtract* when their (undesirable) condition is present. Reward is
clipped to [0, 1] so penalties cannot drive it below zero.

Two entry points share this logic:
* :func:`evaluate` — a host-side rLLM reward_fn (``gdpval_reward_fn``) that
  grades the deliverable text surfaced on the episode. This is the canonical
  spec; the in-sandbox ``tests/grade.py`` written by :mod:`rllm.data.gdpval_builder`
  is a self-contained copy of :func:`grade_rubric` for the case where the
  deliverable file only exists inside the sandbox.
* :func:`grade_rubric` / :func:`extract_deliverable_text` — reusable helpers.

Judge model resolution (first hit wins), mirroring ``claw_eval``:
  1. ``GDPVAL_JUDGE_MODEL`` / ``GDPVAL_JUDGE_BASE_URL`` / ``GDPVAL_JUDGE_API_KEY``.
  2. ``task.metadata`` keys ``judge_model`` / ``judge_base_url`` / ``judge_api_key``.
  3. The user's rLLM provider config (``~/.rllm/config.json``) via litellm.
"""

from __future__ import annotations

import json
import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path

from rllm.eval.types import EvalOutput, Signal
from rllm.types import Episode, Task

logger = logging.getLogger(__name__)

# Hint injected into the solver's system prompt (picked up by
# rllm.eval.reward_fns._resolver.get_verifier_system_prompt).
SYSTEM_PROMPT = (
    "Produce the requested deliverable as an actual file (xlsx/docx/pptx/pdf/csv/txt), "
    "using the exact filename the task specifies, and save it into the 'output/' directory "
    "of your working directory. Your work is graded by an LLM judge applied to that file's "
    "contents, so only the files you place in 'output/' will be evaluated."
)

JUDGE_SYSTEM_PROMPT = """\
You are an impartial grader for GDPval, a benchmark of professional knowledge work.
You are given ONE rubric criterion and the text contents of the deliverable an AI
produced. Decide whether the deliverable satisfies the criterion.

Judge ONLY the single criterion given. Be strict and literal: a criterion is met
only if the deliverable clearly and verifiably satisfies it. If the deliverable is
missing, empty, or the criterion cannot be verified from its contents, it is NOT met.

Respond with ONLY a JSON object: {"met": 0 or 1, "reasoning": "<one sentence>"}."""

JUDGE_USER_TEMPLATE = """\
## Task prompt (context)
{prompt}

## Rubric criterion (weight {score})
{criterion}

## Deliverable contents
{deliverable}

Is this single criterion satisfied? Respond with JSON: {{"met": 0 or 1, "reasoning": "..."}}"""

_MAX_DELIVERABLE_CHARS = 60000


@dataclass
class Criterion:
    score: float
    text: str
    rubric_item_id: str = ""


@dataclass
class GradeResult:
    reward: float
    earned: float
    total_possible: float
    per_criterion: list[dict] = field(default_factory=list)
    ungraded: bool = False
    reason: str = ""


# --------------------------------------------------------------------------- #
# Rubric parsing + aggregation (pure; unit-testable without a judge)
# --------------------------------------------------------------------------- #


def parse_rubric(rubric_json: str | list) -> list[Criterion]:
    """Parse ``rubric_json`` (str or already-decoded list) into Criteria."""
    items = json.loads(rubric_json) if isinstance(rubric_json, str) else rubric_json
    out: list[Criterion] = []
    for it in items or []:
        try:
            score = float(it["score"])
            text = str(it["criterion"]).strip()
        except (KeyError, TypeError, ValueError):
            continue
        if not text:
            continue
        out.append(Criterion(score=score, text=text, rubric_item_id=str(it.get("rubric_item_id", ""))))
    return out


def aggregate(criteria: list[Criterion], met: list[bool]) -> GradeResult:
    """Weight-normalized aggregation. ``met[i]`` = judge says criterion i holds."""
    total_possible = sum(c.score for c in criteria if c.score > 0)
    earned = sum(c.score for c, m in zip(criteria, met, strict=False) if m)
    if total_possible <= 0:
        return GradeResult(reward=0.0, earned=earned, total_possible=total_possible, ungraded=True, reason="no_positive_weight_criteria")
    reward = max(0.0, min(1.0, earned / total_possible))
    per = [{"score": c.score, "met": bool(m), "criterion": c.text, "rubric_item_id": c.rubric_item_id} for c, m in zip(criteria, met, strict=False)]
    return GradeResult(reward=reward, earned=earned, total_possible=total_possible, per_criterion=per)


def grade_rubric(criteria: list[Criterion], deliverable_text: str, prompt: str, judge_fn) -> GradeResult:
    """Judge every criterion via ``judge_fn(criterion_text, score, deliverable, prompt) -> bool``.

    ``judge_fn`` is injected so the aggregation is testable with a stub and the
    same logic can drive either litellm (host) or the openai client (sandbox).
    A criterion whose judge call fails is treated as NOT met (conservative).
    """
    met: list[bool] = []
    for c in criteria:
        try:
            met.append(bool(judge_fn(c.text, c.score, deliverable_text, prompt)))
        except Exception as e:  # noqa: BLE001
            logger.warning("[gdpval] judge failed for criterion %s: %s", c.rubric_item_id or c.text[:40], e)
            met.append(False)
    return aggregate(criteria, met)


# --------------------------------------------------------------------------- #
# Deliverable text extraction
# --------------------------------------------------------------------------- #


def extract_deliverable_text(path: str | Path, max_chars: int = _MAX_DELIVERABLE_CHARS) -> str:
    """Best-effort extraction of a deliverable file's textual content.

    Supports xlsx/xlsm, docx, pptx, pdf, and text-like files. Optional deps
    (openpyxl / python-docx / python-pptx / pdfplumber) are imported lazily;
    a missing dep degrades to an empty/annotated string rather than raising.
    """
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ""
    ext = p.suffix.lower()
    try:
        if ext in (".xlsx", ".xlsm", ".xls"):
            text = _extract_xlsx(p)
        elif ext == ".docx":
            text = _extract_docx(p)
        elif ext == ".pptx":
            text = _extract_pptx(p)
        elif ext == ".pdf":
            text = _extract_pdf(p)
        elif ext in (".csv", ".tsv", ".txt", ".md", ".json", ".html", ".xml"):
            text = p.read_text(encoding="utf-8", errors="replace")
        else:
            # Unknown binary: try utf-8, else note the type.
            try:
                text = p.read_text(encoding="utf-8")
            except (UnicodeDecodeError, OSError):
                text = f"[unsupported deliverable type: {ext or 'no extension'}]"
    except Exception as e:  # noqa: BLE001
        logger.warning("[gdpval] extraction failed for %s: %s", p.name, e)
        text = f"[extraction error for {p.name}: {e}]"
    header = f"# Deliverable file: {p.name}\n"
    body = text if len(text) <= max_chars else text[:max_chars] + "\n...[truncated]..."
    return header + body


def _extract_xlsx(p: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(p, read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_docx(p: Path) -> str:
    import docx  # python-docx

    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(x for x in parts if x is not None)


def _extract_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_pdf(p: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(p)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts.append(f"## Page {i}")
            parts.append(page.extract_text() or "")
    return "\n".join(parts)


# --------------------------------------------------------------------------- #
# Host-side reward_fn entry point
# --------------------------------------------------------------------------- #


def _find_deliverable_text(task: Task, episode: Episode) -> tuple[str, str]:
    """Locate the deliverable's text. Returns (text, source_description).

    FAIL-CLOSED: the grader only ever reads a deliverable the harness has
    explicitly surfaced — never the agent's transcript or final answer. This
    keeps the judge from grading the *process* instead of the *artifact* (and
    from being gamed by an agent that describes a good deliverable it never
    produced). If nothing is surfaced, returns ("", "none") and the caller
    marks the task ungraded rather than scoring it.

    Resolution order:
      1. ``episode.artifacts['deliverable_text']`` — deliverable text surfaced directly.
      2. ``episode.artifacts['deliverable_path']`` / ``output_files`` — a file
         path (sandbox workdir mount or host) to extract.
    """
    arts = episode.artifacts or {}
    if arts.get("deliverable_text"):
        return str(arts["deliverable_text"]), "artifacts.deliverable_text"

    candidates: list[str] = []
    if arts.get("deliverable_path"):
        candidates.append(str(arts["deliverable_path"]))
    of = arts.get("output_files")
    if isinstance(of, list | tuple):
        candidates.extend(str(x) for x in of)
    elif isinstance(of, str):
        candidates.append(of)
    for c in candidates:
        text = extract_deliverable_text(c)
        if text:
            return text, f"file:{Path(c).name}"

    return "", "none"


def _resolve_judge(task: Task) -> tuple[str, str | None, str | None]:
    """Return (model, base_url, api_key); mirrors claw_eval's resolution."""
    env_model = os.environ.get("GDPVAL_JUDGE_MODEL")
    if env_model:
        return env_model, os.environ.get("GDPVAL_JUDGE_BASE_URL"), os.environ.get("GDPVAL_JUDGE_API_KEY")

    meta_model = task.metadata.get("judge_model")
    if meta_model:
        return meta_model, task.metadata.get("judge_base_url"), task.metadata.get("judge_api_key")

    try:
        from rllm.eval.config import get_provider_info, load_config

        cfg = load_config()
        if cfg.provider == "custom":
            return cfg.model, cfg.base_url or None, cfg.api_key or "EMPTY"
        info = get_provider_info(cfg.provider)
        if info and cfg.model:
            prefix = info.litellm_prefix
            model = f"{prefix}/{cfg.model}" if prefix and not cfg.model.startswith(prefix + "/") else cfg.model
            return model, None, cfg.api_key or None
    except Exception:
        logger.debug("[gdpval] could not resolve judge from rLLM config", exc_info=True)
    return "", None, None


def _make_litellm_judge(model: str, base_url: str | None, api_key: str | None):
    import litellm

    def judge_fn(criterion: str, score: float, deliverable: str, prompt: str) -> bool:
        messages = [
            {"role": "system", "content": JUDGE_SYSTEM_PROMPT},
            {"role": "user", "content": JUDGE_USER_TEMPLATE.format(prompt=prompt[:4000], score=score, criterion=criterion, deliverable=deliverable)},
        ]
        # drop_params lets litellm silently drop provider-unsupported params:
        # reasoning models (gpt-5*, o-series) reject temperature=0.0, so it's
        # dropped for them while non-reasoning judges keep deterministic temp=0.
        kwargs: dict = {"model": model, "messages": messages, "temperature": 0.0, "drop_params": True}
        if base_url:
            kwargs["base_url"] = base_url
        if api_key:
            kwargs["api_key"] = api_key
        resp = litellm.completion(**kwargs)
        return _parse_met(resp.choices[0].message.content or "")

    return judge_fn


def _parse_met(text: str) -> bool:
    m = re.search(r"\{.*?\}", text, re.DOTALL)
    if m:
        try:
            return int(json.loads(m.group()).get("met", 0)) == 1
        except (json.JSONDecodeError, ValueError, TypeError):
            pass
    return bool(re.search(r"\b(yes|met|true|1)\b", text, re.IGNORECASE))


def _rubric_as_guidance(criteria: list[Criterion], limit: int = 40) -> str:
    """Render the rubric as bullet guidance for the pairwise judge."""
    if not criteria:
        return "(no rubric provided; judge on overall professional quality)"
    lines = [f"- [weight {int(c.score) if c.score == int(c.score) else c.score}] {c.text}" for c in criteria[:limit]]
    if len(criteria) > limit:
        lines.append(f"- ...(+{len(criteria) - limit} more criteria)")
    return "\n".join(lines)


def evaluate(task: Task, episode: Episode) -> EvalOutput:
    """Host-side ``gdpval_reward_fn``: weighted-rubric LLM-judge grade."""
    rubric_raw = task.metadata.get("rubric_json") or task.metadata.get("rubric")
    criteria = parse_rubric(rubric_raw) if rubric_raw else []
    if not criteria:
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0)],
            metadata={"reason": "no_rubric", "ungraded": True},
        )

    deliverable, source = _find_deliverable_text(task, episode)
    if not deliverable:
        # Fail-closed: no deliverable was surfaced (never grade the transcript).
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0)],
            metadata={"reason": "no_deliverable_surfaced", "ungraded": True},
        )
    prompt = task.metadata.get("prompt") or (task.instruction if isinstance(task.instruction, str) else "") or ""

    model, base_url, api_key = _resolve_judge(task)
    if not model:
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="rubric_score", value=0.0)],
            metadata={"reason": "no_judge_configured", "ungraded": True},
        )

    result = grade_rubric(criteria, deliverable, prompt, _make_litellm_judge(model, base_url, api_key))
    n_met = sum(1 for c in result.per_criterion if c["met"])
    return EvalOutput(
        reward=float(result.reward),
        is_correct=result.reward >= 0.5,
        signals=[Signal(name="rubric_score", value=float(result.reward))],
        metadata={
            "judge_model": model,
            "deliverable_source": source,
            "earned": result.earned,
            "total_possible": result.total_possible,
            "criteria_met": n_met,
            "criteria_total": len(criteria),
            "occupation": task.metadata.get("occupation", ""),
            "sector": task.metadata.get("sector", ""),
        },
    )


# --------------------------------------------------------------------------- #
# Pairwise grader (headline metric) — rubric-informed, vs the gold deliverable
# --------------------------------------------------------------------------- #
#
# This reproduces GDPval's grading standard as closely as a self-hosted judge
# can: the judge compares the model's deliverable against the expert's gold
# deliverable and picks the better one, guided by the task's rubric. Per
# OpenAI's own current guidance (evals.openai.com/gdpval/grading): the standard
# is pairwise human-expert preference, but you may "run an LLM-based judge on
# the rubrics and deliverables to get a rough estimate of model performance."
#
# Output per task: candidate score in {0, 0.5, 1} (loss / tie / win vs. the
# expert), averaged over a position-swapped pair of judge calls to cancel the
# judge's position bias. Averaged across tasks this is the GDPval win-rate.

PAIRWISE_SYSTEM_PROMPT = """\
You are an experienced professional grading two deliverables produced for the same
work request — one by an AI, one by a human expert (you are not told which is which).
Decide which deliverable is the higher-quality professional work product.

Use the provided rubric criteria as your guide, and also weigh professional quality:
correctness, completeness, structure, and fitness for the stated purpose. Be
decisive; only call a genuine tie when the two are truly indistinguishable in quality.

Respond with ONLY a JSON object: {"winner": "A" | "B" | "tie", "reasoning": "<one sentence>"}."""

PAIRWISE_USER_TEMPLATE = """\
## Work request
{prompt}

## Rubric criteria (guidance)
{rubric}

## Deliverable A
{a}

## Deliverable B
{b}

Which deliverable is the better professional work product? Respond with JSON:
{{"winner": "A" | "B" | "tie", "reasoning": "..."}}"""


def _parse_winner(text: str) -> str | None:
    """Return 'A', 'B', or 'tie' from a judge response, or None if unparseable."""
    m = re.search(r"\{.*?\}", text, re.DOTALL)
    if m:
        try:
            w = str(json.loads(m.group()).get("winner", "")).strip().lower()
            if w in ("a", "b", "tie"):
                return w
        except (json.JSONDecodeError, ValueError, TypeError):
            pass
    low = text.lower()
    if "tie" in low:
        return "tie"
    if re.search(r"\bwinner\b[^ab]*b\b", low) or re.search(r"\bb\b.*better", low):
        return "b"
    if re.search(r"\ba\b", low):
        return "a"
    return None


def _make_litellm_pairwise_judge(model: str, base_url: str | None, api_key: str | None):
    import litellm

    def judge(prompt: str, rubric: str, a: str, b: str) -> str | None:
        messages = [
            {"role": "system", "content": PAIRWISE_SYSTEM_PROMPT},
            {"role": "user", "content": PAIRWISE_USER_TEMPLATE.format(prompt=prompt[:4000], rubric=rubric, a=a, b=b)},
        ]
        # drop_params lets litellm silently drop provider-unsupported params:
        # reasoning models (gpt-5*, o-series) reject temperature=0.0, so it's
        # dropped for them while non-reasoning judges keep deterministic temp=0.
        kwargs: dict = {"model": model, "messages": messages, "temperature": 0.0, "drop_params": True}
        if base_url:
            kwargs["base_url"] = base_url
        if api_key:
            kwargs["api_key"] = api_key
        resp = litellm.completion(**kwargs)
        return _parse_winner(resp.choices[0].message.content or "")

    return judge


def score_pairwise(candidate: str, reference: str, prompt: str, rubric: str, judge) -> float | None:
    """Candidate's win score in {0, 0.5, 1}, position-swap debiased.

    ``judge(prompt, rubric, a, b) -> 'a'|'b'|'tie'|None`` is injected for testing.
    Runs twice with candidate in each slot; returns the mean of the two per-call
    candidate scores (win=1, tie=0.5, loss=0). Returns None if both calls fail.
    """

    def _cand_score(winner: str | None, candidate_is: str) -> float | None:
        if winner is None:
            return None
        if winner == "tie":
            return 0.5
        return 1.0 if winner == candidate_is else 0.0

    s1 = _cand_score(judge(prompt, rubric, candidate, reference), "a")  # candidate = A
    s2 = _cand_score(judge(prompt, rubric, reference, candidate), "b")  # candidate = B
    scores = [s for s in (s1, s2) if s is not None]
    if not scores:
        return None
    return sum(scores) / len(scores)


def _find_reference_text(task: Task) -> tuple[str, str]:
    """Locate the gold (expert) deliverable's text. Returns (text, source)."""
    meta = task.metadata or {}
    if meta.get("reference_deliverable_text"):
        return str(meta["reference_deliverable_text"]), "metadata.reference_deliverable_text"
    if meta.get("reference_deliverable_path"):
        text = extract_deliverable_text(meta["reference_deliverable_path"])
        if text:
            return text, "metadata.reference_deliverable_path"
    # Resolve files staged under <task_dir>/reference/ by the builder.
    names = meta.get("reference_deliverables") or []
    if names and task.dataset_dir is not None:
        base = task.dataset_dir / task.sub_dir if task.sub_dir else task.dataset_dir
        for name in names:
            text = extract_deliverable_text(Path(base) / "reference" / name)
            if text:
                return text, f"reference/{name}"
    return "", "none"


def evaluate_pairwise(task: Task, episode: Episode) -> EvalOutput:
    """Headline ``gdpval_pairwise_reward_fn``: rubric-informed pairwise vs. gold.

    reward = candidate win score in {0, 0.5, 1} (loss / tie / win vs. the expert
    deliverable). Mean over tasks = GDPval win-rate.
    """
    candidate, cand_src = _find_deliverable_text(task, episode)
    if not candidate:
        # Fail-closed: no candidate deliverable surfaced (never grade the transcript).
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="pairwise_win", value=0.0)],
            metadata={"reason": "no_deliverable_surfaced", "ungraded": True},
        )
    reference, ref_src = _find_reference_text(task)
    if not reference:
        # No gold deliverable → pairwise is impossible; don't fabricate a score.
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="pairwise_win", value=0.0)],
            metadata={"reason": "no_reference_deliverable", "ungraded": True},
        )

    criteria = parse_rubric(task.metadata.get("rubric_json") or task.metadata.get("rubric") or [])
    rubric = _rubric_as_guidance(criteria)
    prompt = task.metadata.get("prompt") or (task.instruction if isinstance(task.instruction, str) else "") or ""

    model, base_url, api_key = _resolve_judge(task)
    if not model:
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="pairwise_win", value=0.0)],
            metadata={"reason": "no_judge_configured", "ungraded": True},
        )

    score = score_pairwise(candidate, reference, prompt, rubric, _make_litellm_pairwise_judge(model, base_url, api_key))
    if score is None:
        return EvalOutput(
            reward=0.0,
            is_correct=False,
            signals=[Signal(name="pairwise_win", value=0.0)],
            metadata={"reason": "judge_call_failed", "ungraded": True},
        )

    return EvalOutput(
        reward=float(score),
        is_correct=score > 0.5,  # strict win over the expert; tie (0.5) is not a "win"
        signals=[Signal(name="pairwise_win", value=float(score))],
        metadata={
            "judge_model": model,
            "candidate_source": cand_src,
            "reference_source": ref_src,
            "occupation": task.metadata.get("occupation", ""),
            "sector": task.metadata.get("sector", ""),
        },
    )
